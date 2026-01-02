import os
import json
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from jinja2 import Environment, FileSystemLoader

# Configuration
PPTX_FILE = "Iceland_London.pptx"
DATA_FILE = "data.json"
OUTPUT_DIR = "docs"
ASSETS_DIR = "docs/assets"

# Ensure directories exist
os.makedirs(ASSETS_DIR, exist_ok=True)

def load_data():
    with open(DATA_FILE, "r", encoding="utf-8") as f:
        return json.load(f)

def extract_images_and_merge(pptx_path, rich_data):
    prs = pptx.Presentation(pptx_path)
    
    # Create a map of slide_index -> list of image_paths
    slide_images = {}
    
    # We iterate through all slides and extract ALL images
    for i, slide in enumerate(prs.slides):
        slide_images[i] = []
        image_count = 0
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                image_filename = f"image_{i}_{image_count}.{ext}"
                image_path = os.path.join(ASSETS_DIR, image_filename)
                
                # Check if file already exists (optimization)
                if not os.path.exists(image_path):
                    with open(image_path, "wb") as f:
                        f.write(image.blob)
                
                slide_images[i].append(f"assets/{image_filename}")
                image_count += 1
    
    # Merge images into rich_data
    for item in rich_data:
        # If the JSON already explicitly defines an image (e.g. key cover image), keep it (or prepend it)
        existing_images = item.get("images", [])
        
        idx = item.get("image_index")
        extracted = []
        if idx is not None and idx in slide_images:
            extracted = slide_images[idx]
        
        # Strategy: Use existing manually defined images first, then extracted ones
        # This allows us to override PPTX bad images with generated ones if we want
        # But for this user, let's combine: Generated/Manual > Extracted
        item["images"] = existing_images + extracted
            
    return rich_data

def generate_html(data):
    env = Environment(loader=FileSystemLoader("templates"))
    template = env.get_template("index.html")
    
    output_html = template.render(items=data)
    
    with open(os.path.join(OUTPUT_DIR, "index.html"), "w", encoding="utf-8") as f:
        f.write(output_html)

if __name__ == "__main__":
    print(f"Reading rich data from {DATA_FILE}...")
    data = load_data()
    
    print(f"Extracting images from {PPTX_FILE}...")
    merged_data = extract_images_and_merge(PPTX_FILE, data)
    
    print("Generating HTML...")
    generate_html(merged_data)
    print(f"Site generated in {OUTPUT_DIR}/index.html")
