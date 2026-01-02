import pptx

def extract_text(pptx_path):
    prs = pptx.Presentation(pptx_path)
    full_text = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_text.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        full_text.append("\n".join(slide_text))
    
    return "\n\n".join(full_text)

if __name__ == "__main__":
    content = extract_text("Iceland_London.pptx")
    with open("pptx_dump.txt", "w", encoding="utf-8") as f:
        f.write(content)
    print("Dumped to pptx_dump.txt")
