from pptx import Presentation

ppt = Presentation("Iceland_London.pptx")
print(f"Total slides: {len(ppt.slides)}")

for i, slide in enumerate(ppt.slides):
    print(f"\n--- Slide {i+1} ---")
    if slide.shapes.title:
        print(f"Title: {slide.shapes.title.text}")
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            print(f"Text: {shape.text[:50]}...")
        if shape.shape_type == 13: # PICTURE
            print(f"Found Image: {shape.name}")
